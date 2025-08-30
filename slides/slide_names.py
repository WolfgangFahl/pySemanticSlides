#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on 2025-08-29

@author: wf

Slide name formatter CLI.

Provides a small CLI that iterates slides of a .pptx and prints a user-specified
format with variables {page}, {title}, {name}. Built on BaseCmd.
"""

from argparse import ArgumentParser
from pathlib import Path
import re
from typing import Dict, Generator, List, Optional

from basemkit.argparse_action import StoreDictKeyPair
from basemkit.base_cmd import BaseCmd
from pptx import Presentation
from slides.version import Version


class SlideNames:
    """
    Access slide variables from a PowerPoint deck.

    Args:
        path (Path): Path to a .pptx file.

    Attributes:
        prs (Presentation): Loaded presentation instance.
    """

    def __init__(self, path: Path):
        """
        Initialize the slide accessor.

        Args:
            path (Path): Path to a .pptx file.
        """
        self.prs = Presentation(str(path))

    def set_name(self, page: int, new_name: str):
        """
        Set the internal slide name for the given 1-based page.

        Args:
            page (int): 1-based slide index
            new_name (str): New internal name

        Returns:
            tuple: (old_name, new_name)
        """
        if page < 1 or page > len(self.prs.slides):
            raise IndexError(f"page {page} out of range (1..{len(self.prs.slides)})")
        slide = self.prs.slides[page - 1]
        cSld = slide._element.cSld
        old = cSld.get("name")
        cSld.set("name", new_name)
        return old, new_name

    def save(self, out: Path):
        """
        Save the presentation to the given path.

        Args:
            out (Path): Target .pptx file
        """
        self.prs.save(str(out))


    def iter_vars(self) -> Generator[Dict[str, object], None, None]:
        """
        Iterate slide variables.

        Yields:
            dict: A dictionary per slide with:
                - page (int): 1-based slide index.
                - title (str): Slide title text, empty if missing.
                - name (Optional[str]): Internal slide name (cSld@name).
        """
        for idx, slide in enumerate(self.prs.slides, start=1):
            title = slide.shapes.title.text if slide.shapes.title else ""
            # sanitize title
            title= ''.join(title.splitlines())
            title = re.sub(r'\s+', ' ', title)
            title=title.strip()
            name = slide._element.cSld.get("name")
            name = name or "❓"
            yield {"page": idx, "title": title, "name": name}


class SlideNamesCmd(BaseCmd):
    """
    Command-line interface for printing formatted slide lines.

    The format string may reference {page}, {title}, {name}.
    """

    def __init__(self, version: Version):
        """
        Initialize the CLI command.

        Args:
            version (Version): Version metadata object.
        """
        super().__init__(version)
        self.deck: Optional[Path] = None
        self.format_str: str = "{page}\t{name}\t{title}"

    def get_arg_parser(self) -> ArgumentParser:
        """
        Build the argument parser, extending BaseCmd options.

        Returns:
            ArgumentParser: Configured parser with positional deck and -f/--format.
        """
        parser = super().get_arg_parser()
        parser.add_argument("deck", type=Path, help="Path to .pptx file")
        parser.add_argument(
            "--format",
            default=self.format_str,
            help="Format string using {page}, {title}, {name}",
        )
        parser.add_argument(
            "--set",
            action=StoreDictKeyPair,
            metavar="PAGE=NAME[,PAGE=NAME...]",
            help="Set slide names via PAGE=NAME pairs, comma separated",
        )
        parser.add_argument(
            "--out",
            type=Path,
            help="Output .pptx path when using --set (defaults to in-place overwrite)",
        )
        return parser

    def handle_args(self, args) -> bool:
        """
        Handle parsed arguments and execute the command.

        Args:
            args: Parsed argparse namespace.

        Returns:
            bool: True to indicate processing is complete (no further action).
        """
        if super().handle_args(args):
            return True
        self.deck = args.deck
        self.format_str = args.format
        sn = SlideNames(self.deck)
        if args.set:
            out_path = args.out or self.deck
            for page_str, new_name in args.set.items():
                page = int(page_str)
                old, new = sn.set_name(page, new_name)
                print(f"page {page}: [{old}] -> [{new}]")
            sn.save(out_path)
            return True
        for v in sn.iter_vars():
            print(self.format_str.format(**v))
        return True


def main(argv: Optional[List[str]] = None) -> int:
    """
    Entry point.

    Args:
        argv (Optional[List[str]]): Command line arguments or None.

    Returns:
        int: Exit code.
    """
    return SlideNamesCmd.main(Version, argv)


if __name__ == "__main__":
    raise SystemExit(main())
