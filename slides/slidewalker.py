"""
Created on 2022-04-07

@author: wf
"""

import argparse
import csv
import io
import json
import os
import sys
import traceback
import webbrowser
from collections import OrderedDict
from contextlib import redirect_stdout
from io import StringIO
from typing import List

from pptx import Presentation
from tqdm import tqdm

from slides.version import Version


# https://stackoverflow.com/a/70631361/1497139
class YRange:
    """
    an Y Range
    """

    def __init__(self, minY=0, maxY=300):
        self.minY = minY
        self.maxY = maxY

    @staticmethod
    def isIn(yRange, y):
        result = y == 0 or yRange is None or (y >= yRange.minY and y <= yRange.maxY)
        return result


class Slide(object):
    """
    a single slide
    """

    defaultRunDelim = ""

    def __init__(self, ppt, slide, page, pdf_page, runDelim: str = None):
        """
        constructor
        """
        self.ppt = ppt
        self.slide = slide
        self.page = page
        self.pdf_page = pdf_page
        self.name = slide.name
        self.title = None
        if runDelim is None:
            runDelim = Slide.defaultRunDelim
        self.runDelim = runDelim
        # https://stackoverflow.com/a/40821359/1497139
        if slide.shapes.title:
            self.title = slide.shapes.title.text
        if self.title is None:
            self.title = self.name
        pass

    def asDict(self):
        summary = {
            "page": self.page,
            "pdf_page": self.pdf_page,
            "title": self.title,
            "name": self.name,
            "text": self.getText(),
            "notes": self.getNotes(),
        }
        return summary

    def summary(self):
        text = f"{self.page:3d}({self.name}):{self.title}"
        return text

    def getMM(self, emu):
        # https://startbigthinksmall.wordpress.com/2010/01/04/points-inches-and-emus-measuring-units-in-office-open-xml/
        if emu is None:
            return 0
        else:
            return emu.mm

    def getText4Shapes(self, shapes, yRange, runDelim: str = None):
        """
        Get visible text from shapes in a y-range, excluding icon font runs.
        """
        lines = []
        if runDelim is None:
            runDelim = self.runDelim

        for shape in shapes:
            if not shape.has_text_frame:
                continue

            line = ""
            delim = ""
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if any('\ue000' <= c <= '\uf8ff' for c in run.text):
                        continue  # skip icon glyphs
                    line += f"{delim}{run.text}"
                    delim = runDelim

            y = self.getMM(shape.top)
            if y and YRange.isIn(yRange, y) and line.strip():
                lines.append(line.strip())

        return lines


    def getText(self, yRange=None):
        """
        get the text in the given yRange

        Args:
            yRange:

        Return:
            str: the notes for this slide
        """
        text = self.getText4Shapes(self.slide.shapes, yRange, runDelim=self.runDelim)
        return text

    def getNotes(self, yRange=None, useShapes: bool = False) -> str:
        """
        get the notes

        Return:
            str: the notes for this slide
        """
        text = ""
        if self.slide.has_notes_slide:
            notes_slide = self.slide.notes_slide
            if useShapes:
                text = self.getText4Shapes(
                    notes_slide.shapes, yRange, runDelim=self.runDelim
                )
            elif notes_slide.notes_text_frame:
                text = notes_slide.notes_text_frame.text
        return text

    def getLayoutName(self) -> str:
        """
        get the layoutName of this slide
        """
        layoutName = self.slide.slide_layout.name
        return layoutName


class PPT(object):
    """
    PowerPoint Presentation with lecture
    """

    def __init__(self, filepath, pageHeight=297):
        """
        Constructor
        """
        self.filepath = filepath
        self.basename = os.path.basename(filepath)
        self.pageHeight = pageHeight
        if not os.path.isfile(filepath):
            raise Exception("%s does not exist" % filepath)
        self.prs = None
        self.error = None
        self.slides_loaded=False
        self.slides = []

    def summary(self) -> str:
        """
        show a summary of the given lecture
        """
        if self.error:
            summary = f"error: {self.error} at {self.filepath}"
        else:
            if hasattr(self, "lecture"):
                summary = f"{self.title}({len(self.lecture)} lecture)/{self.author}/{self.created}  {self.basename}"
            else:
                summary = f"{self.title}/{self.author}/{self.created}  {self.basename}"
        return summary

    def asDict(self) -> dict:
        """
        convert me to a dict

        Returns:
            dict: summary
        """
        if self.error:
            summary = {"error": str(self.error), "path": self.filepath}
        else:
            summary = {
                "title": self.title,
                "author": self.author,
                "created": self.created,
                "path": self.filepath,
            }
        return summary

    def open(self):
        """
        open my presentation
        """
        try:
            self.prs = Presentation(self.filepath)
            self.author = self.prs.core_properties.author
            self.created = self.prs.core_properties.created
            self.title = self.prs.core_properties.title
        except Exception as ex:
            self.error = ex

    def open_in_office(self):
        """
        open me in the configure office environment
        """
        os.system(f"open {self.filepath}")  # MacOS – adjust for platform

    def getSlides(self, excludeHiddenSlides: bool = False, runDelim: str = None, force: bool = False):
        """
        get my slides

        Args:
            excludeHiddenSlides(bool): if True exclude hidden Slides
            runDelim(str): delimiter for slide text runs
            force(bool): if True, reload slides even if already loaded
        """
        # Return existing slides if already loaded and not forced to reload
        if not force and self.slides_loaded:
            return self.slides
        # Clear existing slides if forcing reload
        if force:
            self.slides = []
        if runDelim is None:
            runDelim = Slide.defaultRunDelim
        if self.prs is None:
            self.open()
        if not self.error:
            page = 0
            pdf_page = 0
            for slide in self.prs.slides:
                page += 1
                if excludeHiddenSlides:
                    if slide._element.get("show") == "0":
                        # slide is hidden → go to next slide
                        continue
                pdf_page += 1
                pptSlide = Slide(
                    self, slide, page=page, pdf_page=pdf_page, runDelim=runDelim
                )
                self.slides.append(pptSlide)
        self.slides_loaded=True
        return self.slides


class PPTSet:
    """
    A set of PowerPoint presentations loaded via a SlideWalker.
    Provides lookup and caching support.
    """

    def __init__(self, slidewalker: "SlideWalker", verbose: bool = False):
        self.slidewalker = slidewalker
        self.verbose = verbose
        self.ppts_by_path: dict[str, PPT] = {}
        self.ppts_by_relpath: dict[str, PPT] = {}

    def load(self, with_progress: bool = False):
        """
        Load presentations using the configured SlideWalker.

        Args:
            with_progress(bool): If True, show a tqdm progress bar.
        """
        ppt_iter = self.slidewalker.yieldPowerPointFiles(verbose=self.verbose)
        iterator = tqdm(ppt_iter, desc="Loading PPTs") if with_progress else ppt_iter
        for ppt in iterator:
            self.ppts_by_path[ppt.filepath] = ppt
            self.ppts_by_relpath[ppt.relpath] = ppt

    def get_ppt(self, path: str, relative: bool = False) -> PPT:
        """
        Retrieve a single presentation by path.

        Args:
            path (str): the relative or absolute file path to the presentation
            relative (bool): if True, lookup by relpath; else, by full path

        Returns:
            PPT: the PowerPoint presentation or None if not loaded
        """
        if relative:
            ppt = self.ppts_by_relpath.get(path)
        else:
            ppt = self.ppts_by_path.get(path)
        return ppt

    def get_slides(self, path: str, relative: bool = False) -> dict[int, Slide]:
        """
        Retrieve slides for a presentation at given path, keyed by page number.

        Args:
            path (str): path to the presentation
            relative (bool): if True, lookup by relpath; else, by full path

        Returns:
            dict[int, Slide]: map from page number to slide
        """
        ppt = self.get_ppt(path, relative=relative)
        slides_by_page: dict[int, Slide] = {}
        if ppt:
            for slide in ppt.getSlides():
                slides_by_page[slide.page] = slide
        return slides_by_page

    def get_slide(self, path: str, page: int, relative: bool = False) -> Slide:
        """
        Get a specific slide by its page number from a presentation.

        Args:
            path (str): path to the presentation
            page (int): 1-based page index
            relative (bool): if True, lookup by relpath; else, by full path

        Returns:
            Slide: the slide object or None if not found
        """
        slides_by_page = self.get_slides(path, relative=relative)
        slide = slides_by_page.get(page)
        return slide


    def as_lod(self) -> List[dict]:
        """
        Return list of dicts representing all presentations.

        Returns:
            List[dict]: list of dicts with presentation metadata
        """
        lod = []
        for ppt in self.ppts_by_path.values():
            record = ppt.asDict()
            lod.append(record)
        return lod


class SlideWalker(object):
    """
    get meta information for all powerpoint presentations in a certain folder
    """

    def __init__(self, rootFolder: str, debug: bool = False):
        """
        Constructor

        Args:
            rootFolder(str): the path to the root folder of the analysis
            debug(bool): if True switch on debugging
        """
        self.rootFolder = rootFolder
        self.debug = debug

    def asCsv(self, listOfDicts: list, fieldNames: list = None) -> str:
        """convert the given list of dicts to CSV
        see https://stackoverflow.com/a/9157370/1497139

        Args:
            listOfDicts(list): the table to convert

        Returns:
            str: the CSV formated result
        """
        output = io.StringIO()
        if fieldNames is None:
            fieldNameSet = set()
            for record in listOfDicts:
                for key in record.keys():
                    fieldNameSet.add(key)
            fieldNames = list(fieldNameSet)
        writer = csv.DictWriter(
            output, fieldnames=fieldNames, quoting=csv.QUOTE_NONNUMERIC
        )
        writer.writeheader()
        for record in listOfDicts:
            writer.writerow(record)
        return output.getvalue()

    def yieldPowerPointFiles(self, verbose: bool = False):
        """
        generate  my power point files

        Args:
            verbose(bool): if True show information about the processing
        """
        pptxFiles = self.findFiles(self.rootFolder, ".pptx")
        if verbose:
            print(f"found {len(pptxFiles)} powerpoint files")
        for pptxFile in pptxFiles:
            if verbose:
                print(f"Extracting data from {pptxFile}")
            ppt = PPT(pptxFile)
            relpath = os.path.relpath(ppt.filepath, self.rootFolder)
            ppt.relpath = relpath
            ppt.open()
            if not ppt.error:
                yield ppt

    def yieldSlides(
        self,
        ppt,
        verbose: bool,
        excludeHiddenSlides: bool = False,
        runDelim: str = None,
        slideDetails: bool = False,
    ):
        """
        yield all slides

        Args:
            verbose(bool): if True print details on stdout
            excludeHiddenSlides(bool): If True hidden lecture will be excluded and also ignored in the page counting
            runDelim(str): the delimiter to use for powerpoint slide text
        """
        ppt.getSlides(excludeHiddenSlides=excludeHiddenSlides, runDelim=runDelim)
        for slide in ppt.slides:
            if verbose and slideDetails:
                print(slide.summary())
            yield slide

    def dumpInfo(
        self,
        outputFormat: str,
        excludeHiddenSlides: bool = False,
        runDelim: str = None,
        slideDetails: bool = False,
    ):
        """
        dump information about the lecture in the given format

        Args:
            outputFormat(str): csv, json or txt
            excludeHiddenSlides(bool): If True hidden lecture will be excluded and also ignored in the page counting
            runDelim(str): the delimiter to use for powerpoint slide text
        """
        info = {}
        csvRecords = []
        verbose = self.debug or outputFormat == "txt"
        for ppt in self.yieldPowerPointFiles(verbose):
            pptSummary = ppt.asDict()
            if verbose:
                print(f"{ppt.summary()}")
            slideSummary = []
            for slide in self.yieldSlides(
                ppt, verbose, excludeHiddenSlides, runDelim, slideDetails=slideDetails
            ):
                slideRecord = slide.asDict()
                csvRecord = OrderedDict()
                csvRecord["basename"] = ppt.basename
                csvRecord["page"] = slideRecord["page"]
                csvRecord["name"] = slideRecord["name"]
                title = "".join(slideRecord["title"].split())
                csvRecord["title"] = title
                csvRecords.append(csvRecord)
                slideSummary.append(slideRecord)
            pptSummary["slides"] = slideSummary
            info[ppt.basename] = pptSummary
        if outputFormat == "json":
            #
            # avoid the windows horror story
            # https://stackoverflow.com/questions/9233027/unicodedecodeerror-charmap-codec-cant-decode-byte-x-in-position-y-character
            # https://stackoverflow.com/a/18337754/1497139
            jsonStr = json.dumps(
                info, indent=2, default=str, ensure_ascii=False
            ).encode("utf8")
            print(jsonStr.decode("utf-8"))
        elif outputFormat == "csv":
            sortedCsvRecords = sorted(
                csvRecords, key=lambda row: (row["basename"], int(row["page"]))
            )
            csvText = self.asCsv(
                sortedCsvRecords, ["basename", "page", "name", "title"]
            )
            print(csvText)
        elif outputFormat == "lod":
            return info

    def dumpInfoToString(self, outputFormat: str, excludeHiddenSlides: bool = True):
        """
        dump information about the presentations in the given format

        Args:
            outputFormat(str): csv, json or txt
            excludeHiddenSlides(bool): If True hidden lecture will be excluded and also ignored in the page counting
        """
        f = StringIO()
        with redirect_stdout(f):
            self.dumpInfo(outputFormat, excludeHiddenSlides=excludeHiddenSlides)
        stdout = f.getvalue()
        return stdout

    def findFiles(self, path: str, ext: str) -> list:
        """
        find Files with the given extension in the given path

        Args:
            path(str): the path to start with
            ext(str): the extension to search for

        Returns:
            list: a list of files found
        """
        foundFiles = []
        for root, _dirs, files in os.walk(path, topdown=False):
            for name in files:
                if name.endswith(ext) and not name.startswith("~$"):
                    filepath = os.path.join(root, name)
                    foundFiles.append(filepath)
        return foundFiles


def main(argv=None):
    """
    main routine
    """
    if argv is None:
        argv = sys.argv
    program_name = os.path.basename(sys.argv[0])
    program_version_message = f"{program_name} (v{Version.version},{Version.updated})"
    try:
        parser = argparse.ArgumentParser(
            description="SlideWalker - get meta information for all powerpoint presentations in a certain folder"
        )
        parser.add_argument(
            "-a",
            "--about",
            help="show about info [default: %(default)s]",
            action="store_true",
        )
        parser.add_argument(
            "-d", "--debug", dest="debug", action="store_true", help="show debug info"
        )
        parser.add_argument(
            "-f",
            "--format",
            default="json",
            help="output format to create: csv,json or txt (default: %(default)s)",
        )
        parser.add_argument(
            "--includeHidden",
            action="store_true",
            help="exclude hidden slides (default: %(default)s)",
        )
        parser.add_argument(
            "--rd",
            "--runDelimiter",
            dest="runDelim",
            help="text run delimiter (default: %(default)s) suggested: ＿↵•",
            default=Slide.defaultRunDelim,
        )
        parser.add_argument("--rootPath", default=".")
        parser.add_argument(
            "-V", "--version", action="version", version=program_version_message
        )
        args = parser.parse_args(argv[1:])
        if args.about:
            print(program_version_message)
            print(f"see {Version.doc_url}")
            webbrowser.open(Version.doc_url)
        else:
            sw = SlideWalker(args.rootPath, args.debug)
            sw.dumpInfo(
                args.format,
                excludeHiddenSlides=not args.includeHidden,
                runDelim=args.runDelim,
            )

    except KeyboardInterrupt:
        ### handle keyboard interrupt ###
        return 1
    except Exception as e:
        indent = len(program_name) * " "
        sys.stderr.write(program_name + ": " + repr(e) + "\n")
        sys.stderr.write(indent + "  for help use --help")
        if args.debug:
            print(traceback.format_exc())
        return 2


if __name__ == "__main__":
    sys.exit(main())
